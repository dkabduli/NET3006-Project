#!/usr/bin/env python3
"""
Generate NET3006A Project Presentation — ML for Network Telemetry
Team: Abdul Rehman, Esam Mukbil, Hashim Kshim, Mazen Alhassan
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── colour palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x1B, 0x1F, 0x3B)   # deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x9F, 0xFD)   # bright cyan-blue
ACCENT_TEAL   = RGBColor(0x00, 0xC9, 0xA7)   # teal green
ACCENT_PURPLE = RGBColor(0x84, 0x5E, 0xC2)   # soft purple
ACCENT_GOLD   = RGBColor(0xF0, 0xB4, 0x29)   # warm gold
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY      = RGBColor(0x99, 0x99, 0x99)
SECTION_BG    = RGBColor(0x0F, 0x14, 0x2E)   # slightly darker navy
TABLE_HEADER  = RGBColor(0x00, 0x7A, 0xC1)
TABLE_ROW_ALT = RGBColor(0x16, 0x1A, 0x35)
ORANGE        = RGBColor(0xFF, 0x8C, 0x00)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── helper functions ────────────────────────────────────────────
def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=0, top=0, width=Inches(0.15), height=None, color=ACCENT_BLUE):
    if height is None:
        height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - height, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = line_spacing
    return txBox


def add_card(slide, left, top, width, height, title, body,
             title_color=WHITE, border_color=ACCENT_BLUE,
             fill_color=RGBColor(0x14, 0x18, 0x32),
             title_size=16, body_size=13, body_color=LIGHT_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.4)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.35),
                title, font_size=title_size, color=title_color, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.55), width - Inches(0.36), height - Inches(0.7),
                body, font_size=body_size, color=body_color, line_spacing=1.3)
    return shape


def add_icon_circle(slide, left, top, size, text, fill_color, font_size=20):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circ.fill.solid()
    circ.fill.fore_color.rgb = fill_color
    circ.line.fill.background()
    add_textbox(slide, left, top + Inches(0.03), size, size - Inches(0.03),
                text, font_size=font_size, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    return circ


def add_image(slide, image_path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(image_path, left, top, **kwargs)


def add_bullet_slide_content(tf, bullets, font_size=18, color=WHITE, bold=False,
                              indent_level=0, font_name='Calibri', spacing=6):
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.level = indent_level
        p.line_spacing = 1.2
        p.space_after = Pt(spacing)
        p.space_before = Pt(1)


def add_presenter_tag(slide, name, color=ACCENT_BLUE):
    add_textbox(slide, Inches(10.2), Inches(6.85), Inches(3), Inches(0.5),
                f"Presenter: {name}", font_size=12, color=color,
                bold=True, alignment=PP_ALIGN.RIGHT)


def add_slide_number(slide, num, total):
    add_textbox(slide, Inches(0.4), Inches(6.9), Inches(1.5), Inches(0.4),
                f"{num} / {total}", font_size=10, color=MID_GRAY,
                alignment=PP_ALIGN.LEFT)


def add_section_header(slide, section_title, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.8), SLIDE_WIDTH, Inches(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x10, 0x15, 0x30)
    shape.line.fill.background()
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.4),
                section_title, font_size=40, color=color, bold=True,
                alignment=PP_ALIGN.CENTER, font_name='Calibri')


def make_title_subtitle(slide, title, subtitle, presenter,
                         title_size=36, subtitle_size=18,
                         title_color=WHITE, subtitle_color=LIGHT_GRAY):
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
                title, font_size=title_size, color=title_color,
                bold=True, font_name='Calibri')
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.12), Inches(11.5), Inches(0.48),
                    subtitle, font_size=subtitle_size, color=subtitle_color,
                    font_name='Calibri')
    add_presenter_tag(slide, presenter)


TOTAL_SLIDES = 23
ASSET_DIR = os.path.join(
    os.path.expanduser("~"),
    ".cursor",
    "projects",
    "Users-ykabduli-Downloads-Third-Year-Second-Semester-NET3006-NET3006-Project",
    "assets",
)
TITLE_VISUAL = os.path.join(ASSET_DIR, "telemetry_title_visual.png")
PIPELINE_VISUAL = os.path.join(ASSET_DIR, "telemetry_pipeline_visual_clean.png")
FUTURE_VISUAL = os.path.join(ASSET_DIR, "future_6g_genai_visual_clean.png")


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
if os.path.exists(TITLE_VISUAL):
    add_image(slide, TITLE_VISUAL, 0, Inches(0.95), width=SLIDE_WIDTH)
add_accent_bar(slide, left=0, top=0, width=SLIDE_WIDTH, height=Inches(0.1), color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)

title_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(0.85), Inches(11.35), Inches(5.45))
title_panel.fill.solid()
title_panel.fill.fore_color.rgb = RGBColor(0x0E, 0x13, 0x2A)
title_panel.fill.transparency = 0.22
title_panel.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.2),
            "Machine Learning for Network Telemetry",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.75), Inches(10.3), Inches(0.7),
            "A Survey of Methods, Applications, and Emerging Trends",
            font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(4.5), Inches(3.3), Inches(4.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_TEAL
shape.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.5),
            "NET3006A — Network Management and Machine Learning",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5),
            "Carleton University  |  Winter 2026  |  Dr. Jie Gao",
            font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5),
            "Abdul Rehman  |  Esam Mukbil  |  Hashim Kshim  |  Mazen Alhassan",
            font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.4),
            "Option 1: Survey / Reading Project  —  Topic 2",
            font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

for i, (label, color) in enumerate([
    ("Detect", ACCENT_BLUE), ("Predict", ACCENT_TEAL), ("Optimize", ACCENT_PURPLE)
]):
    x = Inches(2.35) + Inches(i * 3.0)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.3), Inches(2.2), Inches(0.55))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x32)
    card.line.color.rgb = color
    card.line.width = Pt(1.2)
    add_textbox(slide, x, Inches(6.42), Inches(2.2), Inches(0.3),
                label, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda / Outline
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Presentation Outline", None, "All Members")

items = [
    ("1.", "Introduction & Motivation", "Abdul Rehman"),
    ("2.", "What Network Telemetry Means", "Abdul Rehman"),
    ("3.", "Why Machine Learning Helps", "Esam Mukbil"),
    ("4.", "Anomaly Detection Methods", "Esam Mukbil"),
    ("5.", "Main Lessons from Anomaly Detection", "Mazen Alhassan"),
    ("6.", "Performance Prediction", "Mazen Alhassan"),
    ("7.", "QoS Optimization", "Hashim Kshim"),
    ("8.", "Main Lessons for Performance & QoS", "Hashim Kshim"),
    ("9.", "Telemetry Pipeline", "Mazen Alhassan"),
    ("10.", "Nokia Example", "Mazen Alhassan"),
    ("11.", "Ericsson Example", "Hashim Kshim"),
    ("12.", "Future Trends: 6G & GenAI", "Hashim Kshim"),
    ("13.", "Open Challenges", "Hashim Kshim"),
    ("14.", "Conclusion", "Hashim Kshim"),
]

start_y = Inches(1.6)
for i, (num, title, presenter) in enumerate(items):
    y = start_y + Inches(i * 0.38)
    c = ACCENT_BLUE if "Abdul" in presenter else (ACCENT_TEAL if "Esam" in presenter else (ACCENT_GOLD if "Mazen" in presenter else ACCENT_PURPLE))
    add_textbox(slide, Inches(1.2), y, Inches(0.6), Inches(0.35),
                num, font_size=14, color=c, bold=True)
    add_textbox(slide, Inches(1.8), y, Inches(6), Inches(0.35),
                title, font_size=14, color=WHITE)
    add_textbox(slide, Inches(9.0), y, Inches(3.5), Inches(0.35),
                presenter, font_size=12, color=c, alignment=PP_ALIGN.RIGHT)

# legend
for j, (name, clr) in enumerate([
    ("Abdul Rehman", ACCENT_BLUE), ("Esam Mukbil", ACCENT_TEAL), ("Mazen Alhassan", ACCENT_GOLD), ("Hashim Kshim", ACCENT_PURPLE)
]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.8) + Inches(j * 3.1), Inches(7.05), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    add_textbox(slide, Inches(1.05) + Inches(j * 3.1), Inches(6.95),
                Inches(2.3), Inches(0.35), name, font_size=10, color=clr)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Section Divider: Abdul Rehman
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_bottom_bar(slide, color=ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
            "PART 1", font_size=16, color=ACCENT_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name='Calibri')
add_section_header(slide, "Introduction & ML for Anomaly Detection", ACCENT_BLUE)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Presenter: Abdul Rehman", font_size=20, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — What is Network Telemetry?
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "What is Network Telemetry?",
                    "The network continuously reports what is happening inside it",
                    "Abdul Rehman")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
bullets = [
    "It is the live data that routers, switches, and links send about their health and traffic.",
    "It tells us things like delay, packet loss, throughput, CPU use, and errors.",
    "Modern systems stream this data automatically instead of waiting for slow manual checks.",
    "That gives operators a clearer and faster view of what the network is doing.",
    "Telemetry is the input that ML uses to spot trouble and predict future problems.",
]
add_bullet_slide_content(tf, bullets, font_size=17, color=WHITE, spacing=10)

# right-side visual panel
add_textbox(slide, Inches(7.55), Inches(1.8), Inches(4.8), Inches(0.4),
            "Simple View", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
for x, label, color in [
    (Inches(7.75), "Devices", ACCENT_BLUE),
    (Inches(9.45), "Data", ACCENT_TEAL),
    (Inches(11.15), "ML", ACCENT_PURPLE),
]:
    add_icon_circle(slide, x, Inches(2.45), Inches(1.05), "", color)
    add_textbox(slide, x, Inches(2.78), Inches(1.05), Inches(0.35),
                label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for x in [Inches(8.85), Inches(10.55)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(2.78), Inches(0.45), Inches(0.22))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MID_GRAY
    arrow.line.fill.background()

add_card(
    slide, Inches(7.45), Inches(4.05), Inches(5.45), Inches(2.0),
    "Examples of Telemetry Data",
    "Flow data\nPacket traces\nDevice metrics\nINT path details\nLogs and events",
    title_color=ACCENT_BLUE, border_color=ACCENT_BLUE, body_size=14
)

add_presenter_tag(slide, "Abdul Rehman")
add_slide_number(slide, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Telemetry vs Traditional Monitoring
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "From Traditional Monitoring to Streaming Telemetry",
                    "Why modern networks moved beyond slow polling and basic dashboards",
                    "Abdul Rehman")

add_card(slide, Inches(0.9), Inches(1.9), Inches(5.6), Inches(3.25),
         "Traditional Monitoring",
         "Checks devices every so often\nOften slower and less detailed\nProblems may be noticed after users are affected\nWorks poorly at modern network scale",
         title_color=ACCENT_BLUE, border_color=ACCENT_BLUE, body_size=15)
add_card(slide, Inches(6.85), Inches(1.9), Inches(5.6), Inches(3.25),
         "Streaming Telemetry",
         "Data is pushed continuously\nMore detailed and more real-time\nMakes faster detection and response possible\nBetter fit for ML and automation",
         title_color=ACCENT_TEAL, border_color=ACCENT_TEAL, body_size=15)

arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.95), Inches(3.05), Inches(0.55), Inches(0.3))
arrow.fill.solid()
arrow.fill.fore_color.rgb = MID_GRAY
arrow.line.fill.background()

add_card(slide, Inches(2.5), Inches(5.35), Inches(8.3), Inches(1.05),
         "Main Point",
         "Better telemetry gives ML better input, and better input leads to better decisions.",
         title_color=ACCENT_BLUE, border_color=ACCENT_BLUE, body_size=13)

add_presenter_tag(slide, "Abdul Rehman")
add_slide_number(slide, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Where Telemetry Data Comes From
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Where Telemetry Data Comes From",
                    "Different network signals combine to give a full picture",
                    "Abdul Rehman")

sources = [
    ("Devices", "Routers, switches, and access points report health and counters.", ACCENT_BLUE),
    ("Traffic Flows", "Flow records show who is talking, how much, and how often.", ACCENT_TEAL),
    ("Packets", "Packet traces reveal detailed behavior on the wire.", ACCENT_PURPLE),
    ("Logs & Events", "Logs capture warnings, failures, and configuration changes.", ORANGE),
]
for i, (title, body, color) in enumerate(sources):
    x = Inches(0.9) + Inches((i % 2) * 6.0)
    y = Inches(1.9) + Inches((i // 2) * 2.0)
    add_card(slide, x, y, Inches(5.35), Inches(1.45), title, body,
             title_color=color, border_color=color, body_size=13)

add_card(slide, Inches(2.0), Inches(5.55), Inches(9.3), Inches(1.05),
         "Why this matters",
         "ML works best when it can learn from multiple types of telemetry instead of a single signal.",
         title_color=ACCENT_BLUE, border_color=ACCENT_BLUE, body_size=12)

add_presenter_tag(slide, "Abdul Rehman")
add_slide_number(slide, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Why Telemetry Matters
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Why Telemetry Matters in Modern Networks",
                    "Telemetry supports visibility, reliability, and automation",
                    "Abdul Rehman")

benefits = [
    ("See problems earlier", "Operators can notice trouble before outages spread."),
    ("Understand performance", "Teams can measure delay, loss, and congestion more clearly."),
    ("Support automation", "Telemetry gives ML the evidence it needs to recommend actions."),
    ("Scale to new networks", "5G, cloud, and large systems need better visibility than old tools provide."),
]
for i, (title, body) in enumerate(benefits):
    y = Inches(1.8) + Inches(i * 1.05)
    add_icon_circle(slide, Inches(0.95), y + Inches(0.02), Inches(0.42), str(i + 1), ACCENT_BLUE, font_size=16)
    add_textbox(slide, Inches(1.55), y, Inches(4.2), Inches(0.35),
                title, font_size=18, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(1.55), y + Inches(0.38), Inches(10.8), Inches(0.42),
                body, font_size=14, color=LIGHT_GRAY)

add_card(slide, Inches(7.8), Inches(2.15), Inches(4.6), Inches(2.7),
         "Simple Summary",
         "Telemetry is the network's evidence.\nML uses that evidence to find patterns, explain behavior, and support faster decisions.",
         title_color=ACCENT_TEAL, border_color=ACCENT_TEAL, body_size=15)

add_presenter_tag(slide, "Abdul Rehman")
add_slide_number(slide, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Why ML for Network Telemetry?
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Why Apply ML to Network Telemetry?",
                    "The challenge is not collecting data. It is understanding it fast enough.",
                    "Esam Mukbil")

challenges = [
    ("Too much data", "Large networks create huge streams of measurements every second."),
    ("Problems are hidden", "Bad behavior can come from many small signals happening together."),
    ("Speed matters", "Teams need to detect and react before users feel the problem."),
    ("Networks keep changing", "Traffic, apps, and paths move, so fixed thresholds stop working."),
    ("Automation is needed", "Modern networks need systems that can watch, decide, and respond."),
]

# center the ML label vertically in the circle
ml_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.55), Inches(2.35), Inches(2.0), Inches(2.0))
ml_circle.fill.solid()
ml_circle.fill.fore_color.rgb = ACCENT_BLUE
ml_circle.line.fill.background()
add_textbox(slide, Inches(5.55), Inches(2.95), Inches(2.0), Inches(0.45),
            "ML", font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.95), Inches(4.55), Inches(3.2), Inches(0.35),
            "Why use ML here?", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

card_specs = [
    (Inches(0.9), Inches(1.9), "Too much data", "Large networks create huge telemetry streams every second.", ACCENT_BLUE),
    (Inches(0.9), Inches(4.15), "Problems are hidden", "Small clues across many signals are hard to notice manually.", ACCENT_TEAL),
    (Inches(8.25), Inches(1.9), "Speed matters", "Teams need answers before users feel the problem.", ACCENT_PURPLE),
    (Inches(8.25), Inches(4.15), "Networks change", "Traffic and paths move, so fixed rules become outdated.", ORANGE),
]
for x, y, title, body, color in card_specs:
    add_card(slide, x, y, Inches(3.9), Inches(1.55), title, body,
             title_color=color, border_color=color, body_size=13)

add_card(
    slide, Inches(4.0), Inches(5.75), Inches(5.15), Inches(1.05),
    "What We Cover",
    "Anomalies  |  Prediction  |  QoS  |  Industry  |  Future trends",
    title_color=ACCENT_TEAL, border_color=ACCENT_TEAL, body_size=12
)

add_presenter_tag(slide, "Esam Mukbil", ACCENT_TEAL)
add_slide_number(slide, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — ML Methods for Anomaly Detection
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "ML Methods for Anomaly Detection",
                    "A visual summary of the main model families",
                    "Esam Mukbil")

methods = [
    ("Autoencoders", "Unsupervised", "Learn what normal traffic looks like. If new data looks very different, it may be an anomaly."),
    ("LSTM / GRU", "Deep Learning", "Look at patterns over time. They help catch spikes, drops, or strange sequences in telemetry."),
    ("Isolation Forest", "Unsupervised", "Quickly separates unusual points from common ones. Useful when speed is important."),
    ("One-Class SVM", "Semi-supervised", "Builds a boundary around normal behavior, then flags data outside that boundary."),
    ("GAN-based Models", "Generative", "Generate or compare normal-looking samples so unusual behavior stands out more clearly."),
]

positions = [
    (Inches(1.55), Inches(1.8)), (Inches(5.15), Inches(1.8)), (Inches(8.75), Inches(1.8)),
    (Inches(3.35), Inches(4.0)), (Inches(6.95), Inches(4.0)),
]
badge_colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_BLUE, ACCENT_PURPLE, ORANGE]
for i, (method, category, desc) in enumerate(methods):
    x, y = positions[i]
    add_card(slide, x, y, Inches(3.05), Inches(1.6), method, desc,
             title_color=WHITE, border_color=badge_colors[i], body_size=12)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(1.95), y - Inches(0.18), Inches(0.92), Inches(0.28))
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_colors[i]
    badge.line.fill.background()
    add_textbox(slide, x + Inches(1.95), y - Inches(0.16), Inches(0.92), Inches(0.24),
                category, font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_presenter_tag(slide, "Esam Mukbil", ACCENT_TEAL)
add_slide_number(slide, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Anomaly Detection: Key Findings
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Anomaly Detection: Main Lessons",
                    "Less text, stronger takeaways",
                    "Mazen Alhassan")

findings = [
    "Deep learning often beats older statistical methods when telemetry has many features.",
    "Unsupervised learning is popular because real networks rarely have enough labeled anomaly data.",
    "Hybrid models work well because they mix strong feature learning with fast scoring.",
    "Real-time use is still hard because accurate models can be expensive to run at scale.",
    "Richer telemetry, especially INT, usually improves anomaly detection quality.",
]

summary_cards = [
    ("Deep learning is strong", findings[0], ACCENT_BLUE, Inches(1.3), Inches(1.9)),
    ("Unsupervised is practical", findings[1], ACCENT_TEAL, Inches(7.0), Inches(1.9)),
    ("Hybrid models help", findings[2], ACCENT_PURPLE, Inches(1.3), Inches(3.95)),
    ("Speed is still hard", findings[3], ORANGE, Inches(7.0), Inches(3.95)),
]
for title, body, color, x, y in summary_cards:
    add_card(slide, x, y, Inches(5.0), Inches(1.65), title, body,
             title_color=color, border_color=color, body_size=12)

add_card(
    slide, Inches(2.2), Inches(5.55), Inches(8.9), Inches(1.08),
    "Key Takeaway",
    "Best results often come from models that learn normal behavior first because labeled anomaly data is limited.",
    title_color=ACCENT_BLUE, border_color=ACCENT_BLUE, body_size=12
)

add_presenter_tag(slide, "Mazen Alhassan", ACCENT_GOLD)
add_slide_number(slide, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Section Divider: Esam Mukbil & Mazen Alhassan
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_bottom_bar(slide, color=ACCENT_TEAL)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
            "PART 2", font_size=16, color=ACCENT_TEAL, bold=True,
            alignment=PP_ALIGN.CENTER)
add_section_header(slide, "Methods, Findings & Prediction", ACCENT_TEAL)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Presenters: Esam Mukbil & Mazen Alhassan", font_size=20, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — ML for Performance Prediction
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_TEAL)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "ML for Network Performance Prediction",
                    "Using telemetry to estimate what the network will do next",
                    "Mazen Alhassan")

items_left = [
    ("Time-Series Forecasting", [
        "LSTM and GRU models learn how metrics change over time",
        "They forecast delay, jitter, loss, and throughput",
        "Useful when recent behavior helps predict the near future",
    ]),
    ("Graph Neural Networks (GNNs)", [
        "Represent the network as devices and links in a graph",
        "Good for path-based prediction because topology matters",
        "They capture how one congested link can affect other paths",
    ]),
]

y_pos = Inches(1.7)
for title, bullets in items_left:
    add_textbox(slide, Inches(0.8), y_pos, Inches(5.8), Inches(0.4),
                title, font_size=18, color=ACCENT_TEAL, bold=True)
    y_pos += Inches(0.45)
    for b in bullets:
        add_textbox(slide, Inches(1.1), y_pos, Inches(5.5), Inches(0.4),
                    f"• {b}", font_size=14, color=LIGHT_GRAY)
        y_pos += Inches(0.35)
    y_pos += Inches(0.2)

# right column
items_right = [
    ("Ensemble & Hybrid Methods", [
        "Combine different models to improve robustness",
        "Often more accurate than relying on one model alone",
        "Can adapt better when traffic patterns shift",
    ]),
    ("Attention Mechanisms", [
        "Transformer-style models look at longer patterns in the data",
        "They help identify which signals matter most",
        "Strong option for large, multi-feature telemetry streams",
    ]),
]

y_pos = Inches(1.7)
for title, bullets in items_right:
    add_textbox(slide, Inches(7.0), y_pos, Inches(5.8), Inches(0.4),
                title, font_size=18, color=ACCENT_TEAL, bold=True)
    y_pos += Inches(0.45)
    for b in bullets:
        add_textbox(slide, Inches(7.3), y_pos, Inches(5.5), Inches(0.4),
                    f"• {b}", font_size=14, color=LIGHT_GRAY)
        y_pos += Inches(0.35)
    y_pos += Inches(0.2)

add_presenter_tag(slide, "Mazen Alhassan", ACCENT_GOLD)
add_slide_number(slide, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — ML for QoS Optimization
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_TEAL)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "ML for QoS Optimization",
                    "Using telemetry and predictions to keep service quality high",
                    "Esam Mukbil")

qos_methods = [
    ("Reinforcement Learning", "An agent tries actions and learns which routing or scheduling choices give better QoS."),
    ("Deep Q-Networks", "Use the current network state to choose actions like rerouting, balancing load, or assigning bandwidth."),
    ("Multi-Objective Optimization", "Helps balance trade-offs such as low delay, good throughput, and fairness at the same time."),
    ("Federated Learning", "Lets different domains train shared models without sending all raw telemetry to one place."),
]

for i, (method, desc) in enumerate(qos_methods):
    y = Inches(1.7) + Inches(i * 1.3)
    # accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.8), y, Inches(0.08), Inches(1.0))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_TEAL
    line.line.fill.background()
    add_textbox(slide, Inches(1.1), y, Inches(11.5), Inches(0.35),
                method, font_size=18, color=ACCENT_TEAL, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.38), Inches(11.5), Inches(0.7),
                desc, font_size=14, color=LIGHT_GRAY)

add_presenter_tag(slide, "Esam Mukbil", ACCENT_TEAL)
add_slide_number(slide, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Performance & QoS Key Findings
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_TEAL)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Performance & QoS: Main Lessons",
                    "What stood out most across the literature",
                    "Esam Mukbil")

findings = [
    ("Topology matters",
     "GNNs usually beat plain time-series models when the shape of the network affects performance."),
    ("RL is powerful but risky",
     "Reinforcement learning looks strong in simulations, but deployment is harder because bad actions can hurt live traffic."),
    ("Good data matters most",
     "Missing or noisy telemetry can hurt prediction quality more than model choice."),
    ("Proactive beats reactive",
     "Predicting trouble early is better than waiting for thresholds to fail after users are already affected."),
]

for i, (title, desc) in enumerate(findings):
    y = Inches(1.7) + Inches(i * 1.3)
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.75), y + Inches(0.02), Inches(0.38), Inches(0.38))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = ACCENT_TEAL
    num_shape.line.fill.background()
    add_textbox(slide, Inches(0.75), y + Inches(0.04), Inches(0.38), Inches(0.35),
                str(i + 1), font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.3), y, Inches(11.5), Inches(0.35),
                title, font_size=17, color=ACCENT_TEAL, bold=True)
    add_textbox(slide, Inches(1.3), y + Inches(0.38), Inches(11.5), Inches(0.7),
                desc, font_size=14, color=LIGHT_GRAY)

add_presenter_tag(slide, "Esam Mukbil", ACCENT_TEAL)
add_slide_number(slide, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Telemetry Pipeline Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_TEAL)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "ML-Enhanced Network Telemetry Pipeline",
                    "A cleaner flow from data to action",
                    "Mazen Alhassan")

if os.path.exists(PIPELINE_VISUAL):
    add_image(slide, PIPELINE_VISUAL, Inches(0.7), Inches(1.6), width=Inches(11.95))

stages = [
    ("Collect", "Devices export live network data", ACCENT_BLUE),
    ("Store", "Platforms organize and retain telemetry", ACCENT_TEAL),
    ("Analyze", "ML models detect patterns and predict issues", ACCENT_PURPLE),
    ("Act", "Operators or automation respond faster", ORANGE),
]
for i, (title, desc, color) in enumerate(stages):
    x = Inches(0.8) + Inches(i * 3.05)
    add_card(slide, x, Inches(5.0), Inches(2.7), Inches(1.15), title, desc,
             title_color=color, border_color=color, body_size=11)

# bottom note
add_textbox(slide, Inches(0.9), Inches(6.35), Inches(11.4), Inches(0.45),
            "The pipeline matters because good telemetry architecture makes reliable ML and faster response possible.",
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_presenter_tag(slide, "Mazen Alhassan", ACCENT_GOLD)
add_slide_number(slide, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Section Divider: Mazen Alhassan & Hashim Kshim
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_bottom_bar(slide, color=ACCENT_PURPLE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
            "PART 3", font_size=16, color=ACCENT_GOLD, bold=True,
            alignment=PP_ALIGN.CENTER)
add_section_header(slide, "Pipeline, Industry & Conclusion", ACCENT_GOLD)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Presenters: Mazen Alhassan & Hashim Kshim", font_size=20, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Nokia's Telemetry Solutions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_PURPLE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Industry Example: Nokia",
                    "Showing what ML-driven telemetry looks like in practice",
                    "Mazen Alhassan")

add_card(slide, Inches(0.85), Inches(1.9), Inches(3.8), Inches(1.6),
         "Collect", "Gather live broadband telemetry from many devices.",
         title_color=ACCENT_GOLD, border_color=ACCENT_GOLD)
add_card(slide, Inches(4.75), Inches(1.9), Inches(3.8), Inches(1.6),
         "Analyze", "Use ML to detect anomalies and estimate risk.",
         title_color=ACCENT_TEAL, border_color=ACCENT_TEAL)
add_card(slide, Inches(8.65), Inches(1.9), Inches(3.8), Inches(1.6),
         "Act", "Support faster troubleshooting and predictive maintenance.",
         title_color=ORANGE, border_color=ORANGE)

for x in [Inches(4.2), Inches(8.1)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(2.52), Inches(0.35), Inches(0.22))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MID_GRAY
    arrow.line.fill.background()

add_card(
    slide, Inches(1.45), Inches(4.15), Inches(10.4), Inches(1.7),
    "What Nokia Is Trying To Do",
    "Predict failures early  |  find unusual behavior automatically  |  estimate customer experience  |  forecast capacity needs  |  speed up troubleshooting",
    title_color=ACCENT_GOLD, border_color=ACCENT_GOLD, body_size=14
)

add_presenter_tag(slide, "Mazen Alhassan", ACCENT_GOLD)
add_slide_number(slide, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — Ericsson's AI-Driven Telemetry
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_PURPLE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Industry Example: Ericsson",
                    "Combining better data organization with automation",
                    "Hashim Kshim")

add_card(
    slide, Inches(0.85), Inches(1.75), Inches(5.6), Inches(2.75),
    "AI-Ready Data Mesh",
    "Ericsson treats telemetry as a shared data product, not just raw logs.\n\n"
    "Teams manage domain pipelines with common standards, which improves data quality and makes telemetry easier for ML teams to use.",
    title_color=ACCENT_PURPLE, border_color=ACCENT_PURPLE, body_size=13
)
add_card(
    slide, Inches(6.85), Inches(1.75), Inches(5.6), Inches(2.75),
    "Transport Automation Controller (TAC)",
    "Uses AI/ML to monitor transport performance, detect KPI anomalies, and support proactive capacity planning.\n\n"
    "Direction: closed-loop automation through detect, diagnose, and act workflows.",
    title_color=ACCENT_PURPLE, border_color=ACCENT_PURPLE, body_size=13
)
add_card(
    slide, Inches(1.45), Inches(4.95), Inches(10.4), Inches(1.45),
    "Key Industry Insight",
    "Ericsson shows that good AI outcomes depend on strong telemetry and data foundations. "
    "Better data pipelines make network automation and prediction more realistic at scale.",
    title_color=ACCENT_PURPLE, border_color=ACCENT_PURPLE, body_size=13
)

add_presenter_tag(slide, "Hashim Kshim", ACCENT_PURPLE)
add_slide_number(slide, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — Emerging Trends: 6G & GenAI
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_PURPLE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Emerging Trends: 6G & GenAI",
                    "Two trend groups that shape the future of telemetry",
                    "Hashim Kshim")

trends = [
    ("6G & Next-Gen Telemetry", ACCENT_TEAL, [
        "6G will need even faster decisions and smarter automation",
        "Digital twins can simulate a network before changes are made",
        "Semantic telemetry tries to send only the most useful information",
        "AI may become a built-in network function instead of an extra tool",
    ]),
    ("Generative AI & LLMs for Telemetry", ORANGE, [
        "LLMs may let engineers ask telemetry questions in plain language",
        "GenAI can create synthetic data to help train models",
        "It can also help summarize incidents and reports",
        "Large foundation models may be adapted for different network tasks",
    ]),
]

if os.path.exists(FUTURE_VISUAL):
    add_image(slide, FUTURE_VISUAL, Inches(0.7), Inches(1.55), width=Inches(11.9))
add_card(slide, Inches(0.9), Inches(5.25), Inches(5.65), Inches(1.3),
         trends[0][0], "Edge intelligence, digital twins, semantic telemetry, and AI-native design.",
         title_color=ACCENT_TEAL, border_color=ACCENT_TEAL, body_size=12)
add_card(slide, Inches(6.75), Inches(5.25), Inches(5.65), Inches(1.3),
         trends[1][0], "Natural-language analytics, synthetic data, incident summaries, and foundation models.",
         title_color=ORANGE, border_color=ORANGE, body_size=12)

add_presenter_tag(slide, "Hashim Kshim", ACCENT_PURPLE)
add_slide_number(slide, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — Open Challenges
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_PURPLE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "Open Challenges",
                    "Important problems that still need better solutions",
                    "Hashim Kshim")

challenges = [
    ("Not enough labeled data", "Real networks rarely provide clean examples of every failure or attack."),
    ("Scale and speed", "Models must handle massive telemetry streams quickly enough to matter."),
    ("Explainability", "Operators need to know why a model raised an alert before they trust it."),
    ("Concept drift", "Normal traffic changes over time, so models can become outdated."),
    ("Generalization", "A model that works well in one network may fail in another."),
    ("Privacy and security", "Telemetry can contain sensitive information that must be protected."),
]

for i, (title, desc) in enumerate(challenges):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.8) + col * Inches(6.3)
    y = Inches(1.7) + row * Inches(1.7)
    # card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, Inches(5.8), Inches(1.45))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x32)
    card.line.color.rgb = ACCENT_PURPLE if col == 0 else ORANGE
    card.line.width = Pt(1.2)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.4), Inches(0.35),
                title, font_size=16, color=ACCENT_PURPLE if col == 0 else ORANGE, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.4), Inches(0.75),
                desc, font_size=12, color=LIGHT_GRAY)

add_presenter_tag(slide, "Hashim Kshim", ACCENT_PURPLE)
add_slide_number(slide, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 21 — Conclusion & Key Takeaways
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, left=0, top=0, width=SLIDE_WIDTH, height=Inches(0.1), color=ACCENT_TEAL)
add_bottom_bar(slide, color=ACCENT_BLUE)
make_title_subtitle(slide, "Conclusion & Key Takeaways", "A cleaner final summary", "Hashim Kshim")

takeaways = [
    ("Telemetry gives the raw picture",
     "It tells us what the network is doing right now through streams of measurements and events."),
    ("ML helps turn data into decisions",
     "It can find anomalies, predict future performance, and help improve QoS."),
    ("Industry is already using these ideas",
     "Companies such as Nokia and Ericsson are building more automated, ML-driven telemetry systems."),
    ("The field is still growing",
     "6G, digital twins, and generative AI may make telemetry smarter and easier to use."),
    ("The hardest part is deployment",
     "Real-world success still depends on trustworthy models, good data, and systems that scale."),
]

for i, (title, desc) in enumerate(takeaways):
    y = Inches(1.45) + Inches(i * 1.08)
    # numbered circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.75), y + Inches(0.02), Inches(0.42), Inches(0.42))
    circ.fill.solid()
    colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ORANGE, RGBColor(0xE0, 0x40, 0x40)]
    circ.fill.fore_color.rgb = colors[i]
    circ.line.fill.background()
    add_textbox(slide, Inches(0.75), y + Inches(0.05), Inches(0.42), Inches(0.38),
                str(i + 1), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.4), y, Inches(11.5), Inches(0.35),
                title, font_size=17, color=colors[i], bold=True)
    add_textbox(slide, Inches(1.4), y + Inches(0.38), Inches(11.5), Inches(0.6),
                desc, font_size=14, color=LIGHT_GRAY)

add_presenter_tag(slide, "Hashim Kshim", ACCENT_PURPLE)
add_slide_number(slide, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 22 — Thank You & Q&A
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_accent_bar(slide, left=0, top=0, width=SLIDE_WIDTH, height=Inches(0.1), color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
            "Thank You!", font_size=52, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.7),
            "Questions & Discussion", font_size=30, color=ACCENT_TEAL,
            alignment=PP_ALIGN.CENTER)

# decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5),
            "Abdul Rehman  |  Esam Mukbil  |  Hashim Kshim  |  Mazen Alhassan",
            font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
            "NET3006A — Machine Learning for Network Telemetry",
            font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 23 — References
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, color=ACCENT_BLUE)
add_bottom_bar(slide, color=ACCENT_TEAL)
make_title_subtitle(slide, "References", None, "All Members")

references = [
    "[1]  Nokia, \"Modern Broadband Network Telemetry,\" Nokia White Paper, 2024.",
    "[2]  Ericsson, \"From data mess to AI-ready data mesh,\" Ericsson Blog, 2024.",
    "[3]  Ericsson, \"Visualizing network performance: Transport Automation Controller with AI/ML at the helm,\" Ericsson, 2024.",
    "[4]  Z. Yuan et al., \"iTeleScope: Softwarized Network Middle-Box for Real-Time Video Telemetry and Classification,\" IEEE/ACM Trans. Netw., 2023.",
    "[5]  A. Dogra et al., \"6G Network Architecture: QoS Paradigms and Data Lifecycle Management for Next-Generation Networks,\" IEEE Commun. Surveys Tuts., 2024.",
    "[6]  S. Falkner et al., \"Mobile Network Data Synthesis with Generative AI: Challenges and Solutions,\" IEEE Network, 2024.",
    "[7]  M. Boban et al., \"Autonomous network operations: from reactive management to intent-driven optimization,\" Ericsson Technology Review, 2024.",
    "[8]  Ericsson, \"Four ways generative AI is set to transform the telecom industry,\" Ericsson Blog, 2024.",
    "[9]  ATIS, \"Advancing Generative AI Implementation in Telecommunications Networks,\" ATIS White Paper, 2024.",
    "[10] R. Boutaba et al., \"A comprehensive survey on machine learning for networking,\" ACM Comput. Surv., vol. 51, no. 5, 2018.",
    "[11] P. Mishra et al., \"A detailed investigation and analysis of using ML techniques for intrusion detection,\" IEEE Commun. Surveys Tuts., 2019.",
    "[12] F. Tang et al., \"A survey on machine learning for traffic classification,\" IEEE Trans. on Network and Service Management, 2022.",
]

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8))
tf = txBox.text_frame
tf.word_wrap = True

for i, ref in enumerate(references):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = ref
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Calibri'
    p.space_after = Pt(4)
    p.space_before = Pt(1)

add_slide_number(slide, 23, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "NET3006A_ML_Network_Telemetry_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
